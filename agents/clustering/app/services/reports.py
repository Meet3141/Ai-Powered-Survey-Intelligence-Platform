from __future__ import annotations

import os
from datetime import datetime
from pathlib import Path
from typing import Any

import pandas as pd
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import Image, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle, PageBreak
from pptx import Presentation
from pptx.util import Inches, Pt

from app.core.config import settings


# ============================================================
# HELPERS
# ============================================================

def _ensure_output_dir(output_dir: str | Path | None = None) -> Path:
    path = Path(output_dir or settings.output_dir)
    path.mkdir(parents=True, exist_ok=True)
    return path


def _filter_valid_communities(communities: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """
    Strict filter — discard any community that:
    - has size == 0
    - has no representative_students
    - has no top_interests AND no top_skills
    """
    valid = []
    for c in communities:
        size = c.get("size", 0)
        students = c.get("representative_students", [])
        interests = c.get("top_interests", [])
        skills = c.get("top_skills", [])
        if size < 1:
            continue
        if not students:
            continue
        if not interests and not skills:
            continue
        valid.append(c)
    return valid


def _summary_frame(metrics: dict[str, Any]) -> pd.DataFrame:
    communities = metrics.get("_communities", [])
    rows = [
        {"Metric": "Total Students", "Value": metrics.get("total_students", 0)},
        {"Metric": "Active Communities", "Value": len(communities)},
        {"Metric": "Top Interest", "Value": (metrics.get("interest_distribution") or [{}])[0].get("label", "N/A")},
        {"Metric": "Top Skill", "Value": (metrics.get("skill_distribution") or [{}])[0].get("label", "N/A")},
        {"Metric": "Top Career Goal", "Value": (metrics.get("career_goal_distribution") or [{}])[0].get("label", "N/A")},
    ]
    return pd.DataFrame(rows)


def _get_executive_summary_data(
    metrics: dict[str, Any],
    communities: list[dict[str, Any]],
) -> dict[str, Any]:
    valid_communities = _filter_valid_communities(communities)

    total_students = str(metrics.get("total_students", 0))
    total_communities = str(len(valid_communities))

    largest_community = "N/A"
    largest_size = "0"
    if valid_communities:
        sorted_comms = sorted(valid_communities, key=lambda c: c.get("size", 0), reverse=True)
        largest_community = sorted_comms[0].get("name", "N/A")
        largest_size = str(sorted_comms[0].get("size", 0))

    top_dept = (metrics.get("department_distribution") or [{}])[:1]
    top_dept_name = top_dept[0].get("label", "N/A") if top_dept else "N/A"

    top_interest = (metrics.get("interest_distribution") or [{}])[:1]
    top_interest_name = top_interest[0].get("label", "N/A") if top_interest else "N/A"

    top_skill = (metrics.get("skill_distribution") or [{}])[:1]
    top_skill_name = top_skill[0].get("label", "N/A") if top_skill else "N/A"

    top_goal = (metrics.get("career_goal_distribution") or [{}])[:1]
    top_goal_name = top_goal[0].get("label", "N/A") if top_goal else "N/A"

    recommendations = metrics.get("faculty_recommendations", [])
    top_recommendation = recommendations[0] if recommendations else "No recommendations generated yet."

    return {
        "total_students": total_students,
        "total_communities": total_communities,
        "largest_community": largest_community,
        "largest_size": largest_size,
        "top_dept": top_dept_name.title() if top_dept_name != "N/A" else "N/A",
        "top_interest": top_interest_name.title() if top_interest_name != "N/A" else "N/A",
        "top_skill": top_skill_name.title() if top_skill_name != "N/A" else "N/A",
        "top_goal": top_goal_name.title() if top_goal_name != "N/A" else "N/A",
        "top_recommendation": top_recommendation,
        "recommendations": recommendations,
    }


# ============================================================
# EXCEL REPORT
# ============================================================

def create_excel_report(
    metrics: dict[str, Any],
    communities: list[dict[str, Any]],
    similarity_rows: list[dict[str, Any]],
    recommendations: list[dict[str, Any]],
    chart_paths: dict[str, str],
    output_dir: str | Path | None = None,
) -> str:
    path = _ensure_output_dir(output_dir)
    output_file = path / "faculty_report.xlsx"
    valid_communities = _filter_valid_communities(communities)

    with pd.ExcelWriter(output_file, engine="openpyxl") as writer:
        # Inject communities into metrics for summary frame
        metrics_copy = dict(metrics)
        metrics_copy["_communities"] = valid_communities
        _summary_frame(metrics_copy).to_excel(writer, index=False, sheet_name="Executive Summary")
        pd.DataFrame(valid_communities).to_excel(writer, index=False, sheet_name="Communities")
        pd.DataFrame(similarity_rows).to_excel(writer, index=False, sheet_name="Similarity")
        pd.DataFrame(recommendations).to_excel(writer, index=False, sheet_name="Recommendations")
        pd.DataFrame([{k: str(v) for k, v in metrics.items() if not k.startswith("_")}]).to_excel(
            writer, index=False, sheet_name="Analytics Raw"
        )
        pd.DataFrame([{name: file_path for name, file_path in chart_paths.items()}]).to_excel(
            writer, index=False, sheet_name="Charts"
        )

    return str(output_file)


# ============================================================
# PDF REPORT
# ============================================================

def create_pdf_report(
    metrics: dict[str, Any],
    communities: list[dict[str, Any]],
    summary_text: str,
    chart_paths: dict[str, str],
    output_dir: str | Path | None = None,
) -> str:
    path = _ensure_output_dir(output_dir)
    output_file = path / "faculty_report.pdf"

    doc = SimpleDocTemplate(str(output_file), pagesize=letter)
    styles = getSampleStyleSheet()

    h1 = ParagraphStyle("H1", parent=styles["Heading1"], fontSize=16, spaceAfter=14, spaceBefore=20)
    h2 = ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, spaceAfter=8, spaceBefore=12)
    body = styles["BodyText"]
    label_style = ParagraphStyle("Label", parent=body, fontName="Helvetica-Bold")

    # Filter communities before building the report
    valid_communities = _filter_valid_communities(communities)
    exec_data = _get_executive_summary_data(metrics, valid_communities)
    story: list[Any] = []

    # ── 1. COVER PAGE ────────────────────────────────────────────────────
    story.append(Spacer(1, 150))
    story.append(Paragraph(f"<font size=28><b>{settings.report_title}</b></font>", styles["Title"]))
    story.append(Spacer(1, 20))
    story.append(Paragraph("<font size=16>STEMValley Community Intelligence Platform</font>", styles["Title"]))
    story.append(Spacer(1, 40))
    story.append(Paragraph(f"<font size=14><b>Institution:</b> STEMValley Academic Institute</font>", styles["Title"]))
    story.append(Paragraph(f"<font size=14><b>Generated:</b> {datetime.now().strftime('%B %d, %Y')}</font>", styles["Title"]))
    story.append(Paragraph(f"<font size=14><b>Academic Year:</b> {datetime.now().year}–{datetime.now().year + 1}</font>", styles["Title"]))
    story.append(PageBreak())

    # ── 2. EXECUTIVE SUMMARY ─────────────────────────────────────────────
    story.append(Paragraph("2. EXECUTIVE SUMMARY", h1))

    exec_table_data = [
        ["Metric", "Value"],
        ["Total Students Analyzed", exec_data["total_students"]],
        ["Total Communities Discovered", exec_data["total_communities"]],
        ["Largest Community", f"{exec_data['largest_community']} ({exec_data['largest_size']} students)"],
        ["Top Department", exec_data["top_dept"]],
        ["Top Student Interest", exec_data["top_interest"]],
        ["Top Student Skill", exec_data["top_skill"]],
        ["Top Career Goal", exec_data["top_goal"]],
        ["Top Recommendation", exec_data["top_recommendation"]],
    ]
    exec_table = Table(exec_table_data, colWidths=[220, 300])
    exec_table.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e293b")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 11),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f8fafc"), colors.white]),
        ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
        ("FONTSIZE", (0, 1), (-1, -1), 10),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#e2e8f0")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 10),
        ("RIGHTPADDING", (0, 0), (-1, -1), 10),
        ("TOPPADDING", (0, 0), (-1, -1), 8),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 8),
    ]))
    story.append(exec_table)
    story.append(Spacer(1, 16))
    story.append(Paragraph("<b>AI-Generated Summary:</b>", body))
    story.append(Paragraph(summary_text, body))
    story.append(PageBreak())

    # ── 3. KEY FINDINGS ──────────────────────────────────────────────────
    story.append(Paragraph("3. KEY FINDINGS", h1))
    findings = [
        f"{exec_data['total_students']} students were analyzed in this cohort.",
        f"{exec_data['total_communities']} distinct student communities were discovered.",
        f"The largest community is <b>{exec_data['largest_community']}</b> with {exec_data['largest_size']} members.",
        f"<b>{exec_data['top_dept']}</b> is the most represented department.",
        f"<b>{exec_data['top_skill']}</b> is the most prevalent technical skill.",
        f"<b>{exec_data['top_interest']}</b> is the most dominant student interest.",
        f"{len(metrics.get('top_collaborations', []))} high-confidence collaboration pairs identified.",
        "Strong interdisciplinary interaction detected." if len(metrics.get("department_distribution", [])) > 1
        else "Department-specific clustering detected.",
    ]
    for finding in findings:
        story.append(Paragraph(f"• {finding}", body))
    story.append(PageBreak())

    # ── 4. COMMUNITY INSIGHTS ────────────────────────────────────────────
    story.append(Paragraph("4. COMMUNITY INSIGHTS", h1))

    for community in valid_communities:
        name = community.get("name", "Unknown Community")
        size = community.get("size", 0)
        health = community.get("community_health_score", 0)
        description = community.get("description", "")
        top_interests = community.get("top_interests", [])
        top_skills = community.get("top_skills", [])
        reps = community.get("representative_students", [])
        faculty_actions = community.get("faculty_actions", community.get("faculty_recommendations", []))

        # Health score color
        if health >= 70:
            health_color = colors.HexColor("#10b981")  # green
            health_label = "Excellent"
        elif health >= 45:
            health_color = colors.HexColor("#f59e0b")  # amber
            health_label = "Good"
        else:
            health_color = colors.HexColor("#f43f5e")  # red
            health_label = "Needs Attention"

        story.append(Paragraph(f"<b>{name}</b>", h2))

        # Community summary table
        comm_table_data = [
            ["Field", "Details"],
            ["Description", description or "N/A"],
            ["Member Count", f"{size} students"],
            ["Community Health Score", f"{health}/100 — {health_label}"],
            ["Top Interests", ", ".join(top_interests[:4]) if top_interests else "N/A"],
            ["Top Skills", ", ".join(top_skills[:4]) if top_skills else "N/A"],
            ["Representative Students", ", ".join(reps[:3]) if reps else "N/A"],
        ]
        comm_table = Table(comm_table_data, colWidths=[160, 360])
        comm_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#334155")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.HexColor("#f1f5f9"), colors.white]),
            ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cbd5e1")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.append(comm_table)
        story.append(Spacer(1, 10))

        # Faculty Actions
        story.append(Paragraph("<b>Faculty Actions:</b>", body))
        for action in (faculty_actions or [])[:5]:
            story.append(Paragraph(f"  ▸ {action}", body))
        story.append(PageBreak())

    # ── 5. DEPARTMENT ANALYSIS ───────────────────────────────────────────
    story.append(Paragraph("5. DEPARTMENT ANALYSIS", h1))
    if chart_paths.get("department_distribution"):
        story.append(Image(chart_paths["department_distribution"], width=6 * 72, height=3.5 * 72))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Largest Department:</b> {exec_data['top_dept']}", body))
    story.append(Spacer(1, 6))
    story.append(Paragraph("<b>Department Insights:</b>", body))
    dept_insight = (
        f"The {exec_data['top_dept']} department leads student representation. "
        + (
            "A wide range of departments are participating, showcasing a diverse academic ecosystem."
            if len(metrics.get("department_distribution", [])) > 2
            else "Participation is highly concentrated, highlighting core departmental strengths."
        )
    )
    story.append(Paragraph(dept_insight, body))
    story.append(PageBreak())

    # ── 6. INTEREST ANALYSIS ─────────────────────────────────────────────
    story.append(Paragraph("6. INTEREST ANALYSIS", h1))
    if chart_paths.get("interest_distribution"):
        story.append(Image(chart_paths["interest_distribution"], width=6 * 72, height=3.5 * 72))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Top Interest:</b> {exec_data['top_interest']}", body))
    story.append(Spacer(1, 6))
    story.append(Paragraph("<b>Student Trends & Emerging Interests:</b>", body))
    story.append(Paragraph(
        f"Student focus is heavily oriented toward {exec_data['top_interest']}, reflecting current industry demand. "
        f"Faculty should align coursework and extracurriculars to support these interest clusters.",
        body,
    ))
    story.append(PageBreak())

    # ── 7. SKILL ANALYSIS ────────────────────────────────────────────────
    story.append(Paragraph("7. SKILL ANALYSIS", h1))
    if chart_paths.get("top_skills"):
        story.append(Image(chart_paths["top_skills"], width=6 * 72, height=3.5 * 72))
    story.append(Spacer(1, 12))
    story.append(Paragraph(f"<b>Top Skill:</b> {exec_data['top_skill']}", body))
    story.append(Spacer(1, 6))
    story.append(Paragraph("<b>Industry Readiness & Skill Gaps:</b>", body))
    story.append(Paragraph(
        f"The most prevalent technical capability is {exec_data['top_skill']}. "
        f"Students exhibit strong technical skills matching modern job descriptions. "
        f"Faculty are encouraged to integrate cross-disciplinary soft-skills training to bridge communication gaps.",
        body,
    ))
    story.append(PageBreak())

    # ── 8. STUDENT COLLABORATION ANALYSIS ───────────────────────────────
    story.append(Paragraph("8. STUDENT COLLABORATION ANALYSIS", h1))
    top_collabs = metrics.get("top_collaborations", [])
    if top_collabs:
        for collab in top_collabs[:10]:
            story.append(Paragraph(f"<b>{collab['student_a']} ↔ {collab['student_b']}</b>", h2))
            story.append(Paragraph(f"<b>Similarity Score:</b> {collab['similarity_score']}%", body))
            story.append(Paragraph(
                "<b>Shared Interests:</b> " + (", ".join(collab.get("shared_interests", [])) or "None"), body
            ))
            story.append(Paragraph(
                "<b>Shared Skills:</b> " + (", ".join(collab.get("shared_skills", [])) or "None"), body
            ))
            story.append(Paragraph(f"<b>Collaboration Reason:</b> {collab.get('reason', '')}", body))
            story.append(Spacer(1, 12))
    else:
        story.append(Paragraph("No collaboration data available.", body))
    story.append(PageBreak())

    # ── 9. FACULTY RECOMMENDATIONS ───────────────────────────────────────
    story.append(Paragraph("9. FACULTY RECOMMENDATIONS", h1))
    story.append(Paragraph("<b>Top Priority Actions for Faculty:</b>", body))
    story.append(Spacer(1, 6))
    for rec in exec_data.get("recommendations", []):
        story.append(Paragraph(f"• {rec}", body))
    story.append(PageBreak())

    # ── 10. FUTURE OPPORTUNITIES ─────────────────────────────────────────
    story.append(Paragraph("10. FUTURE OPPORTUNITIES", h1))
    opps = [
        "Interdisciplinary Innovation Hub connecting all communities.",
        "Industry Partnership Programme with top technology employers.",
        "Student Leadership Development Communities across departments.",
    ]
    for c in valid_communities[:5]:
        name = c.get("name", "")
        if "AI" in name or "Machine Learning" in name:
            opps.append("AI Research Programme with GPU compute infrastructure.")
        elif "Entrepreneurship" in name:
            opps.append("Campus Startup Accelerator with seed funding.")
        elif "Cyber" in name:
            opps.append("Cybersecurity Operations Centre for student-run defence practice.")
        elif "Robotics" in name or "IoT" in name:
            opps.append("Robotics & Hardware Innovation Lab with 3D printing and prototyping.")
    opps = list(dict.fromkeys(opps))  # deduplicate
    for opp in opps:
        story.append(Paragraph(f"• {opp}", body))
    story.append(PageBreak())

    # ── 11. APPENDIX ─────────────────────────────────────────────────────
    story.append(Paragraph("11. APPENDIX", h1))
    story.append(Paragraph("<b>Dataset Statistics:</b>", h2))
    story.append(Paragraph(f"• Total Students: {exec_data['total_students']}", body))
    story.append(Paragraph(f"• Total Valid Communities: {exec_data['total_communities']}", body))
    story.append(Paragraph(f"• Clustering Method: {metrics.get('clustering_method', 'N/A')}", body))
    story.append(Paragraph(f"• Final Cluster Count (K): {metrics.get('cluster_count', 'N/A')}", body))
    story.append(Paragraph(
        f"• Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", body
    ))
    story.append(PageBreak())

    # ── 12. METHODOLOGY ──────────────────────────────────────────────────
    story.append(Paragraph("12. METHODOLOGY", h1))
    methods = [
        ("Data Collection", "Survey-based student profiling via STEMValley chatbot interface."),
        ("Data Cleaning", "Agent 2 SurveyClean AI — LLM-assisted normalization + canonical tech-term dictionary with fuzzy matching."),
        ("Embedding Model", "all-MiniLM-L6-v2 — sentence-level semantic embeddings normalized to unit sphere."),
        ("Community Detection", "KMeans with dataset-aware k-bounds (3–5 for ≤50 students) + micro-cluster merging."),
        ("Community Naming", "Canonical taxonomy mapping raw signals to 10 predefined STEM community themes."),
        ("Student Matching", "Cosine similarity on embedding vectors with top-3 match extraction."),
        ("Report Generation", "Agent 3 Community Intelligence Platform — automated PDF, Excel, and PPTX output."),
    ]
    for method_name, method_desc in methods:
        story.append(Paragraph(f"<b>{method_name}:</b> {method_desc}", body))

    doc.build(story)
    return str(output_file)


# ============================================================
# PPTX REPORT
# ============================================================

def create_pptx_report(
    metrics: dict[str, Any],
    communities: list[dict[str, Any]],
    summary_text: str,
    chart_paths: dict[str, str],
    output_dir: str | Path | None = None,
) -> str:
    path = _ensure_output_dir(output_dir)
    output_file = path / "faculty_report.pptx"

    valid_communities = _filter_valid_communities(communities)
    presentation = Presentation()
    exec_data = _get_executive_summary_data(metrics, valid_communities)

    # 1. Cover Slide
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = settings.report_title
    title_slide.placeholders[1].text = (
        f"STEMValley Community Intelligence Platform\n"
        f"Generated: {datetime.now().strftime('%B %d, %Y')}\n"
        f"Academic Year: {datetime.now().year}–{datetime.now().year + 1}"
    )

    # 2. Executive Summary Slide
    exec_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    exec_slide.shapes.title.text = "2. Executive Summary"
    box = exec_slide.shapes.placeholders[1].text_frame
    box.text = f"Total Students: {exec_data['total_students']}"
    for text in [
        f"Total Communities: {exec_data['total_communities']}",
        f"Largest Community: {exec_data['largest_community']} ({exec_data['largest_size']} members)",
        f"Top Department: {exec_data['top_dept']}",
        f"Top Interest: {exec_data['top_interest']}",
        f"Top Skill: {exec_data['top_skill']}",
        f"Top Career Goal: {exec_data['top_goal']}",
        f"Top Recommendation: {exec_data['top_recommendation']}",
    ]:
        p = box.add_paragraph()
        p.text = text
        p.level = 0

    # 3. Key Findings Slide
    find_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    find_slide.shapes.title.text = "3. Key Findings"
    box = find_slide.shapes.placeholders[1].text_frame
    findings = [
        f"{exec_data['total_students']} students analyzed",
        f"{exec_data['total_communities']} distinct communities discovered",
        f"{exec_data['largest_community']} is the largest community",
        f"{exec_data['top_dept']} is the dominant department",
        f"{exec_data['top_skill']} is the most common skill",
        f"{exec_data['top_interest']} is the top student interest",
        f"{len(metrics.get('top_collaborations', []))} high-confidence collaboration pairs identified",
    ]
    box.text = findings[0]
    for text in findings[1:]:
        p = box.add_paragraph()
        p.text = text
        p.level = 0

    # 4. One Slide Per Community
    for community in valid_communities:
        name = community.get("name", "Unknown")
        size = community.get("size", 0)
        health = community.get("community_health_score", 0)
        description = community.get("description", "")
        top_interests = community.get("top_interests", [])
        top_skills = community.get("top_skills", [])
        reps = community.get("representative_students", [])
        actions = community.get("faculty_actions", community.get("faculty_recommendations", []))

        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = f"Community: {name}"
        box = slide.shapes.placeholders[1].text_frame

        box.text = f"Description: {description}"
        for label, val in [
            ("Members", f"{size} students"),
            ("Health Score", f"{health}/100"),
            ("Top Interests", ", ".join(top_interests[:3]) if top_interests else "N/A"),
            ("Top Skills", ", ".join(top_skills[:3]) if top_skills else "N/A"),
            ("Representative Students", ", ".join(reps[:3]) if reps else "N/A"),
            ("Faculty Actions", " | ".join((actions or [])[:3])),
        ]:
            p = box.add_paragraph()
            p.text = f"{label}: {val}"
            p.level = 0

    # 5. Department Analysis Slide
    dept_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    dept_slide.shapes.title.text = "5. Department Analysis"
    if chart_paths.get("department_distribution"):
        dept_slide.shapes.add_picture(chart_paths["department_distribution"], Inches(0.5), Inches(1.5), width=Inches(9))

    # 6. Interest Analysis Slide
    int_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    int_slide.shapes.title.text = "6. Interest Analysis"
    if chart_paths.get("interest_distribution"):
        int_slide.shapes.add_picture(chart_paths["interest_distribution"], Inches(0.5), Inches(1.5), width=Inches(9))

    # 7. Skill Analysis Slide
    skill_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    skill_slide.shapes.title.text = "7. Skill Analysis"
    if chart_paths.get("top_skills"):
        skill_slide.shapes.add_picture(chart_paths["top_skills"], Inches(0.5), Inches(1.5), width=Inches(9))

    # 8. Collaboration Slides
    top_collabs = metrics.get("top_collaborations", [])
    if top_collabs:
        for i in range(0, min(10, len(top_collabs)), 3):
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "8. Student Collaboration Analysis"
            box = slide.shapes.placeholders[1].text_frame
            box.text = ""
            for collab in top_collabs[i : i + 3]:
                p = box.add_paragraph()
                p.text = f"{collab['student_a']} ↔ {collab['student_b']} (Score: {collab['similarity_score']}%)"
                p.level = 0
                p = box.add_paragraph()
                p.text = f"Shared: {', '.join(collab.get('shared_interests', []) + collab.get('shared_skills', []))}"
                p.level = 1
                p = box.add_paragraph()
                p.text = f"Reason: {collab.get('reason', '')}"
                p.level = 1

    # 9. Faculty Recommendations Slide
    rec_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    rec_slide.shapes.title.text = "9. Faculty Recommendations"
    box = rec_slide.shapes.placeholders[1].text_frame
    recs = exec_data.get("recommendations", [])
    if recs:
        box.text = recs[0]
        for text in recs[1:]:
            p = box.add_paragraph()
            p.text = text
            p.level = 0

    # 10. Future Opportunities Slide
    opp_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    opp_slide.shapes.title.text = "10. Future Opportunities"
    box = opp_slide.shapes.placeholders[1].text_frame
    opps = [
        "Interdisciplinary Innovation Hub",
        "Industry Partnership Programme",
        "Student Leadership Communities",
    ]
    for c in valid_communities[:5]:
        n = c.get("name", "")
        if "AI" in n or "Machine" in n:
            opps.append("AI Research Programme")
        elif "Entrepreneur" in n:
            opps.append("Campus Startup Accelerator")
    opps = list(dict.fromkeys(opps))
    box.text = opps[0]
    for text in opps[1:]:
        p = box.add_paragraph()
        p.text = text
        p.level = 0

    # 11. Appendix Slide
    app_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    app_slide.shapes.title.text = "11. Appendix — Dataset Statistics"
    box = app_slide.shapes.placeholders[1].text_frame
    box.text = f"Total Students: {exec_data['total_students']}"
    for text in [
        f"Total Communities: {exec_data['total_communities']}",
        f"Clustering Method: {metrics.get('clustering_method', 'N/A')}",
        f"Final Cluster Count (K): {metrics.get('cluster_count', 'N/A')}",
        f"Report Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
    ]:
        p = box.add_paragraph()
        p.text = text
        p.level = 0

    # 12. Methodology Slide
    meth_slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    meth_slide.shapes.title.text = "12. Methodology"
    box = meth_slide.shapes.placeholders[1].text_frame
    box.text = "Data Collection: Survey-based student profiling"
    for text in [
        "Data Cleaning: LLM + canonical tech-term normalization dictionary",
        "Embedding Model: all-MiniLM-L6-v2",
        "Community Detection: KMeans with dataset-aware k-bounds (3–5 for ≤50 students)",
        "Community Naming: Canonical taxonomy (10 STEM themes)",
        "Student Matching: Cosine Similarity",
        "Report Generation: Agent 3 Community Intelligence Platform",
    ]:
        p = box.add_paragraph()
        p.text = text
        p.level = 0

    presentation.save(output_file)
    return str(output_file)


# ============================================================
# ENTRY POINT
# ============================================================

def build_report_artifacts(
    metrics: dict[str, Any],
    communities: list[dict[str, Any]],
    similarity_rows: list[dict[str, Any]],
    recommendations: list[dict[str, Any]],
    chart_paths: dict[str, str],
    summary_text: str,
    output_dir: str | Path | None = None,
) -> dict[str, str]:
    pdf_path = create_pdf_report(metrics, communities, summary_text, chart_paths, output_dir)
    excel_path = create_excel_report(
        metrics, communities, similarity_rows, recommendations, chart_paths, output_dir
    )
    pptx_path = create_pptx_report(metrics, communities, summary_text, chart_paths, output_dir)
    return {"pdf": pdf_path, "excel": excel_path, "pptx": pptx_path}
